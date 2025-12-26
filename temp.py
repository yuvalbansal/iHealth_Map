# ====================================================================================
#  Bharat_iHealthMap.py
#  Dark Theme – Bharat iHealthMAP (Central Lab style, full corrected version)
# ====================================================================================

from __future__ import annotations

from typing import Optional, List, Tuple

import io
import re

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.io as pio
import plotly.graph_objects as go
import streamlit as st
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE



# --- NEW: optional PPTX support -----------------------------------------------------
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_AUTO_SIZE

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False



# ---------------------------------------------------------------------
#  Plotly dark template + Streamlit page config
# ---------------------------------------------------------------------
pio.templates.default = "plotly_dark"

st.set_page_config(
    page_title="Bharat_iHealthMap",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------
#  Custom CSS – dark background, big fonts
# ---------------------------------------------------------------------
st.markdown(
    """
    <style>
        /* DARK BACKGROUND */
    [data-testid="stAppViewContainer"] {
    background: radial-gradient(circle at top left,
        #161b3f 0, #050816 40%, #02010c 100%);
    font-size: 15px;
    }

    
   /* force all visible content text to gold */
    [data-testid="stAppViewContainer"] p, 
    li {
    color: #E8E8E8; font-size: 20px !important;}

    [data-testid="stAppViewContainer"] h1 { color: #00FF00 !important; }   /* gold */
    [data-testid="stAppViewContainer"] h2 { color: #EFBF04 !important; }   /* cyan */
    [data-testid="stAppViewContainer"] h3 { color: #00FF00 !important; }   /* white */
    [data-testid="stAppViewContainer"] h4 { color: #cccccc !important; }   /* light grey */





    [data-testid="stHeader"] {
        background: rgba(5, 8, 22, 0.1);
    }

    [data-testid="stSidebar"] {
        background: #050816;
    }

    h1, h2, h3, h4 {
        font-weight: 700 !important;
    }

    .metric-card {
        background: #0b1024;
        border-radius: 14px;
        padding: 14px 18px;
        box-shadow: 0 0 16px rgba(15, 23, 42, 0.8);
    }

    .metric-value {
        font-size: 32px;
        font-weight: 800;
        color: #FFD700;
    }

    .metric-label {
        font-size: 14px;
        color: #FFD700;
    }

    .stDataFrame, .stTable {
        font-size: 15px;
    }

    .small-caption {
        font-size: 13px;
        color: #FFD700;
    }

    /* Sidebar text colors */
    [data-testid="stSidebar"] * {
        color: #FFD700 !important;
    }

    /* Sidebar section headers */
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3,
    [data-testid="stSidebar"] h4,
    [data-testid="stSidebar"] h5,
    [data-testid="stSidebar"] h6 {
        color: #FFD700 !important;
        font-weight: 700 !important;
    }

    /* Top nav – we'll underline the area where the radio sits */
    .top-nav-wrapper {
        border-bottom: 3px solid #FF9933; /* saffron */
        padding-bottom: 4px;
        margin-top: 10px;
        margin-bottom: 12px;
    }

    /* Radio styling used as top menu */
    .top-nav-wrapper .stRadio > label {
        font-size: 0; /* hide text label */
    }
    .top-nav-wrapper .stRadio div[role="radiogroup"] {
        display: flex;
        justify-content: center;
        gap: 40px;
    }
    .top-nav-wrapper .stRadio div[role="radiogroup"] > div {
        margin-right: 0 !important;
    }
    .top-nav-wrapper .stRadio div[role="radiogroup"] label {
        font-size: 40px !important;
        font-weight: 600 !important;
        color: #FFD700 !important;
        cursor: pointer;
        padding-bottom: 4px;
    }
    .top-nav-wrapper .stRadio div[role="radiogroup"] input:checked + div > label {
        color: #FFD700 !important;
        border-bottom: 3px solid #FFFFFF;
    }
    




    </style>
    """,
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------
#  Page keys + labels for navbar
# ---------------------------------------------------------------------
PAGE_KEYS = [
    "Overview",
    "Clinical",
    "Lifestyle",
    "Community",
    "Socioeconomic_Profession",
    "Downloads",
]

PAGE_LABELS = {
    "Overview": "Overview",
    "Clinical": "Clinical",
    "Lifestyle": "Lifestyle",
    "Community": "Community",
    "Socioeconomic_Profession": "Socioeconomic & Profession",
    "Downloads": "Downloads",
}

if "page" not in st.session_state:
    st.session_state.page = "Overview"

# ---------------------------------------------------------------------
#  TITLE SECTION (TRICOLOUR UNDERLINE LOOK)
# ---------------------------------------------------------------------
st.markdown(
    """
    <div style="text-align:center; margin-top:10px;">
        <span style="display:block; font-size:44px; font-weight:900; color:#EFBF04;">
            Bharat_iHealthMap
        </span>
        <span style="display:block; font-size:28px; font-weight:600; color:#E8E8E8;">
            Community Health Dashboard
        </span>
        <div style="height:6px; margin-top:6px;
            background: linear-gradient(to right,
                #FF9933 0%,  /* saffron */
                #FFFFFF 50%, /* white */
                #138808 100% /* green */
            );
            border-radius:4px;">
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------
#  TOP BAR DATA UPLOAD
# ---------------------------------------------------------------------
st.markdown("<div style='height:20px'></div>", unsafe_allow_html=True)

upload_col1, upload_col2, upload_col3 = st.columns([2, 4, 2])

with upload_col2:
    st.markdown(
        """
        <div style="
            text-align:center;
            padding:8px;
            border-radius:10px;
            background:rgba(255,255,255,0.08);
            border:1px solid rgba(255,255,255,0.15);
            backdrop-filter: blur(4px);
        ">
            <span style="font-size:18px; color:#FFFFFF; font-weight:200;">
                📁 Upload Data (.xlsx)
            </span>
        </div>
        """,
        unsafe_allow_html=True,
    )

    uploaded_file = st.file_uploader(
        "Upload Excel file",
        type=["xlsx"],
        label_visibility="collapsed",
        help="Upload an Excel (.xlsx) file for population analytics.",
    )

# ---------------------------------------------------------------------
#  TOP NAVIGATION – internal Streamlit session_state (no URL hacks)
# ---------------------------------------------------------------------
with st.container():
    st.markdown('<div class="top-nav-wrapper">', unsafe_allow_html=True)
    selected_page = st.radio(
        "Main navigation",
        options=PAGE_KEYS,
        index=PAGE_KEYS.index(st.session_state.page)
        if st.session_state.page in PAGE_KEYS
        else 0,
        format_func=lambda key: PAGE_LABELS[key],
        horizontal=True,
        key="topnav_radio",
    )
    st.markdown("</div>", unsafe_allow_html=True)

st.session_state.page = selected_page
page = selected_page  # used later in the big if/elif

# ---------------------------------------------------------------------
#  Stop early if no file yet
# ---------------------------------------------------------------------
if uploaded_file is None:
    st.stop()

# ---------------------------------------------------------------------
#  Cached Excel loader
# ---------------------------------------------------------------------
@st.cache_data(show_spinner=True)
def load_excel(uploaded_file):
    """Safely read Excel using openpyxl if possible."""
    try:
        return pd.read_excel(uploaded_file, engine="openpyxl")
    except TypeError:
        # Fallback for older Pandas/Excel engines
        return pd.read_excel(uploaded_file)


df_raw = load_excel(uploaded_file)
source_name = uploaded_file.name

st.caption(
    f"📄 Data source: **{source_name}** — rows: {len(df_raw):,}, "
    f"columns: {len(df_raw.columns)}"
)

df = df_raw.copy()

# ---------------------------------------------------------------------
#  Make dataframe Arrow-compatible
# ---------------------------------------------------------------------

def _safe_numeric(series: pd.Series) -> pd.Series:
    """
    Convert numeric-looking strings to numbers.
    Leave pure text columns untouched.
    """
    try:
        converted = pd.to_numeric(series, errors="coerce")
        # If conversion creates at least one real number, keep it
        if converted.notna().any():
            return converted
        return series
    except Exception:
        return series


for col in df.columns:
    if df[col].dtype == "object":
        df[col] = _safe_numeric(df[col])

# ---------------------------------------------------------------------
#  Helper functions
# ---------------------------------------------------------------------
def norm(s: str) -> str:
    """Normalize a column name for fuzzy matching."""
    return re.sub(r"[^a-z0-9]+", "", str(s).lower())


def find_col(df: pd.DataFrame, candidates: List[str]) -> Optional[str]:
    """Find first existing column among candidate names (case & punctuation insensitive)."""
    if df is None:
        return None
    lookup = {norm(c): c for c in df.columns}
    for name in candidates:
        key = norm(name)
        if key in lookup:
            return lookup[key]
    # contains-style search
    for c in df.columns:
        if any(norm(tok) in norm(c) for tok in candidates):
            return c
    return None


def as_num(series: pd.Series) -> pd.Series:
    return pd.to_numeric(series, errors="coerce")


def parse_diet(val) -> str:
    """Return 'Vegetarian', 'Non-Vegetarian', or 'Unknown'."""
    s = str(val).strip().lower()
    if not s or s == "nan":
        return "Unknown"
    s = s.replace("-", " ").replace("/", " ")
    if "non" in s and "veg" in s:
        return "Non-Vegetarian"
    if "vegitarian" in s:
        return "Vegetarian"
    if "veg" in s or "vegetarian" in s:
        return "Vegetarian"
    return "Unknown"


def flag_yes(val) -> str:
    """Map free-text yes/no-ish answers to Yes/No/Unknown."""
    s = str(val).strip().lower()
    if not s or s == "nan":
        return "Unknown"
    if any(tok in s for tok in ["yes", "y", "positive", "present"]):
        return "Yes"
    if any(tok in s for tok in ["no", "nil", "none", "absent"]):
        return "No"
    return "Unknown"


def donut_normal_abnormal(
    series: pd.Series,
    label: str,
    thresholds: Tuple[Optional[float], Optional[float]],
) -> px.pie:
    """
    Build a donut chart for Normal vs Abnormal.
    thresholds = (low_ok, high_ok). Values outside are 'Abnormal'.
    If low_ok is None => only high cutoff. If high_ok is None => only low cutoff.
    """
    s = pd.to_numeric(series, errors="coerce")
    s = s[s.notna()]

    if len(s) == 0:
        df_plot = pd.DataFrame({"Category": ["No data"], "Count": [1]})
    else:
        low_ok, high_ok = thresholds
        if low_ok is not None and high_ok is not None:
            normal = s.between(low_ok, high_ok)
        elif low_ok is None and high_ok is not None:
            normal = s <= high_ok
        elif low_ok is not None and high_ok is None:
            normal = s >= low_ok
        else:
            normal = pd.Series(False, index=s.index)

        abnormal = ~normal
        df_plot = pd.DataFrame(
            {
                "Category": ["Normal", "Abnormal"],
                "Count": [int(normal.sum()), int(abnormal.sum())],
            }
        )

    fig = px.pie(
        df_plot,
        values="Count",
        names="Category",
        hole=0.55,
        title=label,
    )
    fig.update_traces(textinfo="percent+label")
    fig.update_layout(showlegend=False, margin=dict(l=0, r=0, t=40, b=0))
    return fig


# ---------------------------------------------------------------------
#  Column mapping
# ---------------------------------------------------------------------
cols: dict[str, Optional[str]] = {}

cols["reg_date"] = find_col(df, ["RegistrationDate", "Reg Date", "Date"])
cols["lab_id"] = find_col(
    df,
    [
        "LabId",
        "LAB ID",
        "Lab ID",
        "College/ School/Community Registration Number",
        "Registration Number",
    ],
)
cols["name"] = find_col(df, ["PatientName", "Name", "Student Name", "Participant Name"])
cols["gender"] = find_col(df, ["Gender", "Sex"])
cols["age"] = find_col(df, ["Age", "Age Years", "Age (Years)", "Age in Years"])

cols["locality"] = find_col(
    df,
    ["ADDRESS", "Address", "Locality", "Area", "School/College", "College/ School/Community"],
)
cols["education"] = find_col(
    df,
    ["EDUCATION/HIGHEST QUALIFICATION", "Education/Highest Qualification", "Education"],
)
cols["occupation"] = find_col(df, ["Occupation", "Profession"])
cols["income"] = find_col(df, ["Income", "Monthly Income", "Income Group"])

cols["glucose"] = find_col(df, ["GLUCOSE, FASTING", "Fasting Glucose", "GLUCOSE FASTING"])
cols["chol"] = find_col(df, ["CHOLESTEROL", "Total Cholesterol"])
cols["creatinine"] = find_col(df, ["CREATININE"])
cols["alt"] = find_col(df, ["ALT (SGPT)", "ALT", "SGPT"])
cols["protein_total"] = find_col(df, ["PROTEIN, TOTAL", "Total Protein"])
cols["albumin"] = find_col(df, ["PROTEIN, Albumin ", "PROTEIN, Albumin", "Albumin"])
cols["globulin"] = find_col(df, ["PROTEIN Globulin", "Globulin"])
cols["agr"] = find_col(df, ["PROTEIN G Ratio", "A/G RATIO"])

cols["weight"] = find_col(df, ["WEIGHT IN KG", "Weight (kg)", "Weight"])
cols["height"] = find_col(df, ["HEIGHT", "Height (cm)", "Height"])
cols["bmi"] = find_col(df, ["BMI"])

cols["waist"] = find_col(df, ["WAIST IN INCHES", "Waist", "Waist (inches)"])
cols["bp_sys"] = find_col(df, ["BP SYSTOLIC", "Systolic BP", "BP Systolic", "Systolic"])
cols["bp_dia"] = find_col(df, ["BP DIASTOLIC", "Diastolic BP", "BP Diastolic", "Diastolic"])

cols["tobacco"] = find_col(df, ["HISTORY OF TOBACCO/ SMOKING", "Smoking", "Tobacco"])
cols["alcohol"] = find_col(df, ["HISTORY OF ALCOHOL/ DRUGS", "Alcohol"])
cols["history"] = find_col(
    df,
    ["History of Diabetes/Hypertension / Heart /Kidney/ Liver disease:", "Medical History"],
)
cols["fam_cancer"] = find_col(df, ["Family history of cancer"])
cols["diet"] = find_col(df, ["DIETARY RECALL", "Diet", "Vegetarian/Non-Vegetarian"])
cols["sleep"] = find_col(df, ["SLEEP PATTERN", "Sleep Pattern"])
cols["wake_time"] = find_col(df, ["WAKE UP AT", "Wake up time"])

# ---------------------------------------------------------------------
#  Basic cleaning & derived columns
# ---------------------------------------------------------------------

# 1️⃣ Age -> numeric helper column __AGE__
age_col = cols["age"]
if age_col:
    df["__AGE__"] = df[age_col].astype(str).str.extract(r"(\d+\.?\d*)")[0]
    df["__AGE__"] = pd.to_numeric(df["__AGE__"], errors="coerce")
else:
    df["__AGE__"] = np.nan

# 2️⃣ BMI compute if missing
if cols["bmi"] is None and cols["height"] and cols["weight"]:
    h_cm = as_num(df[cols["height"]])
    w_kg = as_num(df[cols["weight"]])
    with np.errstate(divide="ignore", invalid="ignore"):
        bmi_calc = w_kg / (h_cm / 100.0) ** 2
    df["__BMI__"] = bmi_calc.round(1)
    cols["bmi"] = "__BMI__"

# 3️⃣ Coerce numeric for lab/vital columns
for key in [
    "bmi",
    "bp_sys",
    "bp_dia",
    "glucose",
    "chol",
    "creatinine",
    "alt",
    "protein_total",
    "albumin",
    "globulin",
    "waist",
    "height",
    "weight",
]:
    c = cols.get(key)
    if c:
        df[c] = as_num(df[c])

# 4️⃣ ID column
ID_COL = cols["lab_id"] or cols["name"] or "__index__"
if ID_COL == "__index__":  # if we had to create a synthetic ID
    df["__index__"] = np.arange(1, len(df) + 1)

# ---------------------------------------------------------------------
#  Sidebar filters – simple, global
# ---------------------------------------------------------------------
st.sidebar.header("🔍 Filters")

# Age slider (single definition – uses __AGE__)
age_vals = df["__AGE__"].dropna()
if len(age_vals) > 0:
    amin = int(age_vals.min())
    amax = int(age_vals.max())
    if amin == amax:
        slider_min = max(0, amin - 1)
        slider_max = amin + 1
    else:
        slider_min = max(0, amin)
        slider_max = max(amax, amin + 1)
    age_min, age_max = st.sidebar.slider(
        "Age range (years)",
        min_value=slider_min,
        max_value=slider_max,
        value=(slider_min, slider_max),
        help="All analytics below will respect this filtered age range.",
    )
else:
    age_min, age_max = 0, 120

# Gender filter
if cols["gender"]:
    gender_vals = sorted(df[cols["gender"]].dropna().astype(str).unique().tolist())
    gender_sel = st.sidebar.multiselect("Gender", options=gender_vals, default=gender_vals)
else:
    gender_sel = []

# Diet
if cols["diet"]:
    all_diet = df[cols["diet"]].map(parse_diet)
    diet_vals = sorted(all_diet.dropna().unique().tolist())
    diet_sel = st.sidebar.multiselect("Diet type", options=diet_vals, default=diet_vals)
else:
    diet_sel = []

# Tobacco + Alcohol
tob_sel: List[str] = []
alc_sel: List[str] = []
if cols["tobacco"]:
    tob_vals = sorted(df[cols["tobacco"]].dropna().map(flag_yes).unique().tolist())
    tob_sel = st.sidebar.multiselect("Tobacco / smoking", options=tob_vals, default=tob_vals)
if cols["alcohol"]:
    alc_vals = sorted(df[cols["alcohol"]].dropna().map(flag_yes).unique().tolist())
    alc_sel = st.sidebar.multiselect("Alcohol / drugs", options=alc_vals, default=alc_vals)

# Apply filters (on full df)
mask = pd.Series(True, index=df.index)

# Age
mask &= df["__AGE__"].between(age_min, age_max) | df["__AGE__"].isna()

# Gender
if cols["gender"] and gender_sel:
    mask &= df[cols["gender"]].astype(str).isin(gender_sel)

# Diet
if cols["diet"] and diet_sel:
    mask &= df[cols["diet"]].map(parse_diet).isin(diet_sel)

# Tobacco
if cols["tobacco"] and tob_sel:
    mask &= df[cols["tobacco"]].map(flag_yes).isin(tob_sel)

# Alcohol
if cols["alcohol"] and alc_sel:
    mask &= df[cols["alcohol"]].map(flag_yes).isin(alc_sel)

filtered = df[mask].copy()
st.sidebar.success(f"Filtered records: {len(filtered):,} / {len(df):,}")

# ---------------------------------------------------------------------
#  Health assessment
# ---------------------------------------------------------------------
def assess_health(row: pd.Series) -> pd.Series:
    diagnosis, prognosis, rx = [], [], []

    bmi_val = row.get(cols["bmi"], np.nan) if cols["bmi"] else np.nan
    sys = row.get(cols["bp_sys"], np.nan) if cols["bp_sys"] else np.nan
    dia = row.get(cols["bp_dia"], np.nan) if cols["bp_dia"] else np.nan
    glu = row.get(cols["glucose"], np.nan) if cols["glucose"] else np.nan
    chol_val = row.get(cols["chol"], np.nan) if cols["chol"] else np.nan
    creat = row.get(cols["creatinine"], np.nan) if cols["creatinine"] else np.nan
    alt_val = row.get(cols["alt"], np.nan) if cols["alt"] else np.nan

    diet_raw = row.get(cols["diet"], None) if cols["diet"] else None
    diet = parse_diet(diet_raw)

    # BP
    if pd.notna(sys) and pd.notna(dia):
        if sys >= 180 or dia >= 120:
            diagnosis.append("Hypertensive crisis")
            prognosis.append("Immediate stroke / organ damage risk.")
            rx.append("Urgent medical evaluation.")
        elif sys >= 140 or dia >= 90:
            diagnosis.append("Hypertension")
            prognosis.append("High risk of CVD and stroke.")
            rx.append("Low-salt diet, weight control, activity; consult physician.")
        elif sys >= 120 or dia >= 80:
            diagnosis.append("Pre-hypertension")
            prognosis.append("Likely to progress without lifestyle change.")
            rx.append("Lifestyle optimisation, regular BP monitoring.")

    # Glucose
    if pd.notna(glu):
        if glu >= 126:
            diagnosis.append("Diabetes (fasting)")
            prognosis.append("Micro/macro-vascular complication risk.")
            rx.append("HbA1c, medical review, diet + activity plan.")
        elif 100 <= glu < 126:
            diagnosis.append("Pre-diabetes (fasting)")
            prognosis.append("High risk of progression to diabetes.")
            rx.append("Weight reduction, 150+ min/wk activity, low refined carbs.")

    # Cholesterol
    if pd.notna(chol_val):
        if chol_val >= 240:
            diagnosis.append("High cholesterol")
            prognosis.append("Atherosclerotic CVD risk.")
            rx.append("Reduce saturated/trans fats, full lipid panel.")
        elif 200 <= chol_val < 240:
            diagnosis.append("Borderline high cholesterol")
            prognosis.append("Moderate CVD risk.")
            rx.append("Dietary modification, recheck in 3–6 months.")

    # BMI
    if pd.notna(bmi_val):
        if bmi_val >= 30:
            diagnosis.append("Obesity")
            prognosis.append("Metabolic syndrome & CVD risk.")
            rx.append("Structured weight management, caloric deficit, exercise.")
        elif 25 <= bmi_val < 30:
            diagnosis.append("Overweight")
            prognosis.append("Increased cardiometabolic risk.")
            rx.append("Increase activity, portion control, diet optimisation.")

    # Renal
    if pd.notna(creat) and creat > 1.2:
        diagnosis.append("Possible renal impairment")
        prognosis.append("Needs eGFR & urine protein assessment.")
        rx.append("Avoid nephrotoxins, check renal profile.")

    # Hepatic
    if pd.notna(alt_val) and alt_val > 40:
        diagnosis.append("Elevated ALT (transaminitis)")
        prognosis.append("Fatty liver / hepatitis / drug effect possible.")
        rx.append("Review alcohol, obesity, hepatotoxic drugs; liver evaluation.")

    # Lifestyle
    tob_flag = flag_yes(row.get(cols["tobacco"])) if cols["tobacco"] else "Unknown"
    alc_flag = flag_yes(row.get(cols["alcohol"])) if cols["alcohol"] else "Unknown"
    if tob_flag == "Yes":
        diagnosis.append("Tobacco exposure")
        prognosis.append("↑ risk of CVD, COPD, cancer.")
        rx.append("Cessation counselling, nicotine replacement as needed.")
    if alc_flag == "Yes":
        diagnosis.append("Alcohol / substance use")
        prognosis.append("Liver, mental health, injury risk.")
        rx.append("Limit intake; consider de-addiction help for heavy use.")

    if diet == "Vegetarian":
        rx.append("Ensure adequate protein: dals, soy, paneer, nuts & seeds.")
    elif diet == "Non-Vegetarian":
        rx.append("Prefer fish & skinless poultry; limit fried & red meat.")

    if not diagnosis:
        status = "Healthy"
        diagnosis.append("No major risk flags detected.")
        prognosis.append("Maintain healthy lifestyle & periodic screening.")
        rx.append("Balanced diet, physical activity, regular check-ups.")
    elif len(diagnosis) == 1:
        status = "At Risk"
    else:
        status = "Needs Attention"

    return pd.Series(
        {
            "Health Status": status,
            "Diagnosis": "; ".join(dict.fromkeys(diagnosis)),
            "Prognosis": "; ".join(dict.fromkeys(prognosis)),
            "Prescription": "; ".join(dict.fromkeys(rx)),
        }
    )


if len(filtered) > 0:
    assessed = filtered.apply(assess_health, axis=1)
    view = pd.concat([filtered.reset_index(drop=True), assessed.reset_index(drop=True)], axis=1)
else:
    view = filtered.copy()

# ---------------------------------------------------------------------
#  OVERVIEW
# ---------------------------------------------------------------------
if page == "Overview":
    st.header("📊 Population Overview")

    if len(view) == 0:
        st.info("No records match the current filters. Adjust filters from the sidebar.")
    else:
        total_all = len(df)
        total_filt = len(view)

        status_counts = (
            view["Health Status"].value_counts()
            .reindex(["Healthy", "At Risk", "Needs Attention"])
            .fillna(0)
            .astype(int)
        )
        healthy = int(status_counts.get("Healthy", 0))
        at_risk = int(status_counts.get("At Risk", 0))
        needs_attn = int(status_counts.get("Needs Attention", 0))

        c1, c2, c3, c4 = st.columns(4)
        for col, label, value in [
            (c1, "Total in dataset", f"{total_all:,}"),
            (c2, "Total (filtered)", f"{total_filt:,}"),
            (c3, "Healthy", f"{healthy:,}"),
            (c4, "Needs attention", f"{needs_attn:,}"),
        ]:
            with col:
                st.markdown(
                    f"""
<div class="metric-card">
  <div class="metric-label">{label}</div>
  <div class="metric-value">{value}</div>
</div>
""",
                    unsafe_allow_html=True,
                )

        st.markdown("### Health status distribution")

        colA, colB = st.columns(2)

        with colA:
            fig = px.pie(
                names=status_counts.index,
                values=status_counts.values,
                hole=0.55,
                title="Health status (filtered population)",
            )
            fig.update_traces(textinfo="percent+label")
            st.plotly_chart(fig, width="stretch")

        with colB:
            if cols["gender"]:
                g_counts = (
                    view[cols["gender"]].astype(str).value_counts().reset_index()
                )
                g_counts.columns = ["Gender", "Count"]
                fig_g = px.bar(
                    g_counts,
                    x="Gender",
                    y="Count",
                    title="Sample distribution by gender",
                    text="Count",
                )
                fig_g.update_traces(textposition="outside")
                st.plotly_chart(fig_g, width="stretch")

        st.markdown("### Key biochemical risk flags (approximate)")

        risk_cards = []
        if cols["glucose"]:
            g = as_num(view[cols["glucose"]])
            risk_cards.append(
                ("High fasting glucose (≥126)", f"{100 * (g >= 126).mean():.1f}%")
            )
        if cols["chol"]:
            cchol = as_num(view[cols["chol"]])
            risk_cards.append(
                ("High cholesterol (≥240)", f"{100 * (cchol >= 240).mean():.1f}%")
            )
        if cols["bmi"]:
            b = as_num(view[cols["bmi"]])
            risk_cards.append(("Obesity (BMI ≥30)", f"{100 * (b >= 30).mean():.1f}%"))

        if risk_cards:
            rc1, rc2, rc3 = st.columns(3)
            cols_rc = [rc1, rc2, rc3]
            for idx, (label, val) in enumerate(risk_cards):
                if idx >= len(cols_rc):
                    break
                with cols_rc[idx]:
                    st.markdown(
                        f"""
<div class="metric-card">
  <div class="metric-label">{label}</div>
  <div class="metric-value">{val}</div>
</div>
""",
                        unsafe_allow_html=True,
                    )

        st.markdown("### Parameter-wise normal vs abnormal (donut style)")

        d1, d2, d3 = st.columns(3)
        if cols["glucose"]:
            with d1:
                fig_g = donut_normal_abnormal(view[cols["glucose"]], "Fasting glucose", (70, 99))
                st.plotly_chart(fig_g, width="stretch")
        if cols["chol"]:
            with d2:
                fig_c = donut_normal_abnormal(view[cols["chol"]], "Total cholesterol", (0, 199))
                st.plotly_chart(fig_c, width="stretch")
        if cols["bmi"]:
            with d3:
                b = as_num(view[cols["bmi"]])
                fig_b = donut_normal_abnormal(b, "BMI (18.5–24.9 normal)", (18.5, 24.9))
                st.plotly_chart(fig_b, width="stretch")

        st.markdown("### Sample records with automated health assessment")
        st.dataframe(view.head(100), width="stretch")
        st.markdown(
            '<div class="small-caption">Only first 100 rows shown for display; all '
            "calculations are done on the full filtered dataset.</div>",
            unsafe_allow_html=True,
        )

# ---------------------------------------------------------------------
#  CLINICAL
# ---------------------------------------------------------------------
elif page == "Clinical":
    st.header("🩺 Clinical indicators & risk stratification")

    if len(view) == 0:
        st.info("No records match the filters.")
    else:
        col1, col2 = st.columns(2)

        # Fasting glucose
        if cols["glucose"]:
            with col1:
                st.subheader("Fasting glucose (mg/dL)")
                fig = px.histogram(
                    view,
                    x=cols["glucose"],
                    nbins=40,
                    title="Fasting glucose distribution",
                )
                fig.update_layout(xaxis_title="Glucose (mg/dL)")
                st.plotly_chart(fig, width="stretch")

                g = as_num(view[cols["glucose"]])
                cat = {
                    "Normal (<100)": (g < 100).mean(),
                    "Prediabetes (100–125)": ((g >= 100) & (g < 126)).mean(),
                    "Diabetes (≥126)": (g >= 126).mean(),
                }
                df_cat = pd.DataFrame(
                    {
                        "Category": list(cat.keys()),
                        "Share (%)": [v * 100 for v in cat.values()],
                    }
                )
                fig_cat = px.bar(
                    df_cat,
                    x="Category",
                    y="Share (%)",
                    title="Glycemic categories (approx.)",
                    text="Share (%)",
                )
                fig_cat.update_traces(texttemplate="%{text:.1f}", textposition="outside")
                fig_cat.update_layout(yaxis_title="Share (%)")
                st.plotly_chart(fig_cat, width="stretch")

        # Cholesterol
        if cols["chol"]:
            with col2:
                st.subheader("Total cholesterol (mg/dL)")
                figc = px.histogram(
                    view,
                    x=cols["chol"],
                    nbins=40,
                    title="Cholesterol distribution",
                )
                figc.update_layout(xaxis_title="Cholesterol (mg/dL)")
                st.plotly_chart(figc, width="stretch")

                cchol = as_num(view[cols["chol"]])
                catc = {
                    "Desirable (<200)": (cchol < 200).mean(),
                    "Borderline (200–239)": ((cchol >= 200) & (cchol < 240)).mean(),
                    "High (≥240)": (cchol >= 240).mean(),
                }
                dfc = pd.DataFrame(
                    {
                        "Category": list(catc.keys()),
                        "Share (%)": [v * 100 for v in catc.values()],
                    }
                )
                figc2 = px.bar(
                    dfc,
                    x="Category",
                    y="Share (%)",
                    title="Cholesterol risk categories",
                    text="Share (%)",
                )
                figc2.update_traces(texttemplate="%{text:.1f}", textposition="outside")
                figc2.update_layout(yaxis_title="Share (%)")
                st.plotly_chart(figc2, width="stretch")

        # Blood pressure – improved 3D density view
        if cols["bp_sys"] and cols["bp_dia"]:
            st.subheader("3D Blood Pressure Density (Systolic × Diastolic)")

            bp_df = view[[cols["bp_sys"], cols["bp_dia"]]].dropna().astype(float)

            # Count frequency of each pair
            bp_counts = (
                bp_df.groupby([cols["bp_sys"], cols["bp_dia"]])
                .size()
                .reset_index(name="Count")
            )

            # Build improved 3D scatter
            fig_bp3d = px.scatter_3d(
                bp_counts,
                x=cols["bp_sys"],
                y=cols["bp_dia"],
                z="Count",
                color="Count",
                size="Count",
                opacity=0.9,
                title="3D BP Density (Systolic × Diastolic)",
                color_continuous_scale=["#138808", "#FFFFFF", "#FF9933"],  
            )

            # FIX: Make chart large and beautiful
            fig_bp3d.update_layout(
                width=1000,
                height=700,
                margin=dict(l=0, r=0, t=80, b=20),

                scene=dict(
                    xaxis_title="Systolic (mmHg)",
                    yaxis_title="Diastolic (mmHg)",
                    zaxis_title="Number of persons",

                    # Keep axes proportional & clean
                    aspectmode="cube",

                    # GRID colors
                    xaxis=dict(showbackground=True, backgroundcolor="black"),
                    yaxis=dict(showbackground=True, backgroundcolor="black"),
                    zaxis=dict(showbackground=True, backgroundcolor="black"),
                ),

                # FIX: Best camera angle
                scene_camera=dict(
                    eye=dict(x=1.8, y=1.8, z=1.4),  
                    center=dict(x=0, y=0, z=0)
                )
            )

            st.plotly_chart(fig_bp3d, width="stretch")



        # BMI
        if cols["bmi"]:
            st.subheader("Body mass index (BMI)")
            fig_bmi = px.histogram(
                view,
                x=cols["bmi"],
                nbins=40,
                title="BMI distribution",
            )
            fig_bmi.update_layout(xaxis_title="BMI (kg/m²)")
            st.plotly_chart(fig_bmi, width="stretch")

# ---------------------------------------------------------------------
#  LIFESTYLE
# ---------------------------------------------------------------------
elif page == "Lifestyle":
    st.header("🍽️ Lifestyle & behaviour patterns")

    if len(view) == 0:
        st.info("No records match the filters.")
    else:
        col1, col2 = st.columns(2)

        if cols["diet"]:
            with col1:
                st.subheader("Diet pattern")
                d_series = view[cols["diet"]].map(parse_diet)
                d_counts = d_series.value_counts()
                fig_d = px.pie(
                    values=d_counts.values,
                    names=d_counts.index,
                    hole=0.55,
                    title="Vegetarian vs non-vegetarian (self-reported)",
                )
                st.plotly_chart(fig_d, width="stretch")

        if cols["sleep"]:
            with col2:
                st.subheader("Sleep pattern (top 8 responses)")
                s_counts = view[cols["sleep"]].astype(str).value_counts().head(8)
                fig_s = px.bar(
                    x=s_counts.index,
                    y=s_counts.values,
                    labels={"x": "Sleep pattern", "y": "Count"},
                    title="Sleep pattern (most common)",
                    text=s_counts.values,
                )
                fig_s.update_traces(textposition="outside")
                st.plotly_chart(fig_s, width="stretch")

        col3, col4 = st.columns(2)

        if cols["tobacco"]:
            with col3:
                st.subheader("Tobacco / smoking")
                tob = view[cols["tobacco"]].map(flag_yes).value_counts()
                fig_t = px.bar(
                    x=tob.index,
                    y=tob.values,
                    labels={"x": "Tobacco history", "y": "Count"},
                    title="Tobacco / smoking history",
                    text=tob.values,
                )
                fig_t.update_traces(textposition="outside")
                st.plotly_chart(fig_t, width="stretch")

        if cols["alcohol"]:
            with col4:
                st.subheader("Alcohol / drugs")
                alc = view[cols["alcohol"]].map(flag_yes).value_counts()
                fig_a = px.bar(
                    x=alc.index,
                    y=alc.values,
                    labels={"x": "Alcohol/drugs", "y": "Count"},
                    title="Alcohol / substance use",
                    text=alc.values,
                )
                fig_a.update_traces(textposition="outside")
                st.plotly_chart(fig_a, width="stretch")

        if cols["tobacco"] and cols["alcohol"]:
            st.subheader("Combined lifestyle risk (Tobacco × Alcohol)")

            combo_df = (
                view.assign(
                    Tobacco=view[cols["tobacco"]].map(flag_yes),
                    Alcohol=view[cols["alcohol"]].map(flag_yes),
                )
                .groupby(["Tobacco", "Alcohol"])
                .size()
                .reset_index(name="Count")
            )

            combo_df["Group"] = (
                combo_df["Tobacco"] + " Tobacco / " + combo_df["Alcohol"] + " Alcohol"
            )
            combo_df["Percentage"] = (
                combo_df["Count"] * 100 / combo_df["Count"].sum()
            ).round(1)

            fig_comb = px.bar(
                combo_df,
                x="Group",
                y="Percentage",
                text="Percentage",
                title="Lifestyle Risk Pattern (Tobacco × Alcohol)",
                color="Group",
                color_discrete_sequence=["#138808", "#FF9933", "#FFFFFF", "#00A3E0"],
            )

            fig_comb.update_layout(
                xaxis_title="Risk behaviour group",
                yaxis_title="Population percentage (%)",
                showlegend=False,
            )
            fig_comb.update_traces(textposition="outside")

            st.plotly_chart(fig_comb, width="stretch")


# ---------------------------------------------------------------------
#  COMMUNITY / LOCALITY
# ---------------------------------------------------------------------
elif page == "Community":
    st.header("🏘️ Community / locality health overview")

    loc_col = cols["locality"]
    if not loc_col:
        st.info("No locality/address column detected in the dataset.")
    elif len(view) == 0:
        st.info("No records match the current filters.")
    else:
        agg = view.copy()
        agg["is_unhealthy"] = (agg["Health Status"] != "Healthy").astype(int)
        if cols["glucose"]:
            agg["high_glu"] = (agg[cols["glucose"]] >= 126).astype(int)
        else:
            agg["high_glu"] = 0
        if cols["bmi"]:
            agg["obese"] = (agg[cols["bmi"]] >= 30).astype(int)
        else:
            agg["obese"] = 0

        grp = (
            agg.groupby(loc_col)
            .agg(
                Population=(ID_COL, "count"),
                Unhealthy_rate=("is_unhealthy", "mean"),
                Diabetes_rate=("high_glu", "mean"),
                Obesity_rate=("obese", "mean"),
            )
            .reset_index()
        )

        grp["Unhealthy_rate"] *= 100
        grp["Diabetes_rate"] *= 100
        grp["Obesity_rate"] *= 100

        # Choose how many localities to display
        max_locs = min(50, len(grp))
        n_loc = st.slider(
            "Number of localities to display",
            min_value=5,
            max_value=max_locs,
            value=min(25, max_locs),
            step=5,
            help="Affects both the population and risk charts below.",
        )

        grp_top = grp.sort_values("Population", ascending=False).head(n_loc)

        st.subheader("Top localities by population (filtered subset)")
        fig_pop = px.bar(
            grp_top,
            x=loc_col,
            y="Population",
            title="Top localities by number of records",
            text="Population",
            color="Population",
            color_continuous_scale=["#138808", "#FFFFFF", "#FF9933"],
        )
        fig_pop.update_traces(textposition="outside")
        fig_pop.update_layout(xaxis_tickangle=-45)
        st.plotly_chart(fig_pop, width="stretch")

        st.subheader("Locality-wise risk indicators")
        fig_loc = px.bar(
            grp_top,
            x=loc_col,
            y=["Unhealthy_rate", "Diabetes_rate", "Obesity_rate"],
            barmode="group",
            title="Community risk comparison (Unhealthy / Diabetes / Obesity)",
        )
        fig_loc.update_layout(
            xaxis_title="Locality",
            yaxis_title="Rate (%)",
            xaxis_tickangle=-65,
        )
        st.plotly_chart(fig_loc, width="stretch")


        st.markdown("### Locality table (top 100 by population)")
        st.dataframe(grp.sort_values("Population", ascending=False).head(100), width="stretch")

# ---------------------------------------------------------------------
#  SOCIOECONOMIC & PROFESSION
# ---------------------------------------------------------------------
elif page == "Socioeconomic_Profession":
    st.header("👥 Socioeconomic & profession-wise patterns")

    if len(view) == 0:
        st.info("No records match the filters.")
    else:
        # Income groups (if present)
        if cols["income"]:
            st.subheader("Income group distribution & health status")
            inc_counts = view[cols["income"]].astype(str).value_counts().head(10)
            fig_inc = px.bar(
                x=inc_counts.index,
                y=inc_counts.values,
                title="Top income categories (filtered)",
                labels={"x": "Income group", "y": "Count"},
                text=inc_counts.values,
            )
            fig_inc.update_traces(textposition="outside")
            fig_inc.update_layout(xaxis_tickangle=-30)
            st.plotly_chart(fig_inc, width="stretch")

            if "Health Status" in view:
                inc_status = (
                    view[[cols["income"], "Health Status"]]
                    .dropna()
                    .groupby([cols["income"], "Health Status"])
                    .size()
                    .reset_index(name="Count")
                )
                fig_inc_h = px.bar(
                    inc_status,
                    x=cols["income"],
                    y="Count",
                    color="Health Status",
                    barmode="group",
                    title="Health status by income group (top categories)",
                )
                fig_inc_h.update_layout(xaxis_tickangle=-30)
                st.plotly_chart(fig_inc_h, width="stretch")

        # Occupation / profession
        if cols["occupation"]:
            st.subheader("Occupation / profession distribution & risk")

            occ_counts = view[cols["occupation"]].astype(str).value_counts().head(15)
            fig_occ = px.bar(
                x=occ_counts.index,
                y=occ_counts.values,
                title="Top professions (filtered)",
                labels={"x": "Profession", "y": "Count"},
                text=occ_counts.values,
            )
            fig_occ.update_traces(textposition="outside")
            fig_occ.update_layout(xaxis_tickangle=-30)
            st.plotly_chart(fig_occ, width="stretch")

            if "Health Status" in view:
                occ_status = (
                    view[[cols["occupation"], "Health Status"]]
                    .dropna()
                    .groupby([cols["occupation"], "Health Status"])
                    .size()
                    .reset_index(name="Count")
                )

                # Keep only top 15 professions by total count
                rank = (
                    occ_status.groupby(cols["occupation"])["Count"]
                    .sum()
                    .sort_values(ascending=False)
                    .head(15)
                    .index
                )
                occ_status_top = occ_status[occ_status[cols["occupation"]].isin(rank)]

                fig_occ2 = px.bar(
                    occ_status_top,
                    x=cols["occupation"],
                    y="Count",
                    color="Health Status",
                    barmode="group",
                    title="Health status distribution by profession (Top 15)",
                    text="Count",
                    color_discrete_sequence=["#138808", "#FF9933", "#00A3E0"],
                )
                fig_occ2.update_layout(
                    xaxis_title="Profession",
                    yaxis_title="Number of persons",
                    xaxis_tickangle=-65,
                )
                fig_occ2.update_traces(textposition="outside")

                st.plotly_chart(fig_occ2, width="stretch")


# ---------------------------------------------------------------------
#  DOWNLOADS – Excel + individual PDFs + PPT SUMMARY (NEW)
# ---------------------------------------------------------------------
# ---------------------------------------------------------------------
#  DOWNLOADS – Excel + individual PDFs + PPT SUMMARY (NEW)
# ---------------------------------------------------------------------
elif page == "Downloads":
    st.header("📥 Downloads & individual reports")

    if len(view) == 0:
        st.info("No records match the current filters; adjust filters first.")
    else:
        # -----------------------------------------------------------------
        #  Excel download
        # -----------------------------------------------------------------
        st.subheader("Download filtered records (.xlsx)")

        def to_excel_bytes(df_: pd.DataFrame) -> bytes:
            out = io.BytesIO()
            with pd.ExcelWriter(out, engine="openpyxl") as xw:
                df_.to_excel(xw, index=False, sheet_name="CentralLab_iHealthMAP")
            out.seek(0)
            return out.getvalue()

        st.download_button(
            label="💾 Download filtered records (.xlsx)",
            data=to_excel_bytes(view),
            file_name="central_lab_ihealthmap_filtered.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

        # -----------------------------------------------------------------
        #  Population-level PPT summary (slide-wise analytics)
        # -----------------------------------------------------------------
        st.markdown("---")
        st.subheader("📊 Generate population PPT summary (slide-wise)")

        from pptx.enum.text import MSO_AUTO_SIZE
        from pptx.util import Pt

        # Session variables so PPT can be generated once and downloaded
        if "ppt_version" not in st.session_state:
            st.session_state.ppt_version = 0
        if "ppt_buffer" not in st.session_state:
            st.session_state.ppt_buffer = None

        if not PPTX_AVAILABLE:
            st.info(
                "Install the python-pptx package to enable PPT download "
                "(e.g. pip install python-pptx) and restart the app."
            )
        else:
            col_p1, col_p2 = st.columns(2)
            with col_p1:
                ppt_location = st.text_input("Location for title slide (City/Village/State)")
            with col_p2:
                ppt_year = st.text_input("Year for title slide")

            def safe_pct(num: float, den: float) -> float:
                return float(num) * 100.0 / float(den) if den not in (0, None) else 0.0

            def set_autofit(tf, maxsize=28, minsize=14):
                """
                Apply auto-fit + shrink-to-fit safely.
                """
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                try:
                    tf.fit_text(max_size=maxsize, min_size=minsize)
                except Exception:
                    pass

            class DesignTheme:
                """
                Centralized theme configuration for the presentation.
                Theme: Clean Medical/Tech
                """
                # Colors (RGB)
                PRIMARY = RGBColor(0, 91, 150)       # Deep Teal/Blue #005b96
                SECONDARY = RGBColor(100, 151, 177)  # Soft Blue #6497b1
                ACCENT = RGBColor(255, 111, 105)     # Coral/Orange #ff6f69
                DARK_BG = RGBColor(3, 37, 76)        # Dark Night Blue #03254c (for Titles)
                LIGHT_BG = RGBColor(255, 255, 255)   # White
                TEXT_MAIN = RGBColor(50, 50, 50)     # Dark Grey
                TEXT_LIGHT = RGBColor(240, 240, 240) # Off-white for dark backgrounds
                
                # Fonts
                HEAD_FONT = "Arial"
                BODY_FONT = "Arial"
                
                # Plotly Template Name
                PLOTLY_TEMPLATE = "plotly_white"
                
                # Color Sequence for Charts
                COLOR_SEQUENCE = ["#005b96", "#6497b1", "#ff6f69", "#03254c", "#b3cde0"]

            def configure_plotly_theme():
                """Returns a layout template for Plotly charts to match PPT style."""
                return dict(
                    font_family=DesignTheme.BODY_FONT,
                    font_color="#333333",
                    font_size=16,
                    title_font_family=DesignTheme.HEAD_FONT,
                    title_font_size=24,
                    colorway=DesignTheme.COLOR_SEQUENCE,
                    plot_bgcolor="rgba(0,0,0,0)",
                    paper_bgcolor="rgba(0,0,0,0)",
                    xaxis=dict(
                        showgrid=False, 
                        linecolor="#333",
                        title_font_size=20,
                        tickfont_size=16
                    ),
                    yaxis=dict(
                        showgrid=True, 
                        gridcolor="#eee", 
                        linecolor="#333",
                        title_font_size=20,
                        tickfont_size=16
                    ),
                    legend=dict(
                        font_size=18
                    ),
                )

            def donut_normal_abnormal(series, label, thresholds):
                s = pd.to_numeric(series, errors="coerce")
                s = s[s.notna()]

                if len(s) == 0:
                    df_plot = pd.DataFrame({"Category": ["No data"], "Count": [1]})
                else:
                    low_ok, high_ok = thresholds
                    if low_ok is not None and high_ok is not None:
                        normal = s.between(low_ok, high_ok)
                    elif low_ok is None and high_ok is not None:
                        normal = s <= high_ok
                    elif low_ok is not None and high_ok is None:
                        normal = s >= low_ok
                    else:
                        normal = pd.Series(False, index=s.index)

                    abnormal = ~normal
                    df_plot = pd.DataFrame(
                        {
                            "Category": ["Normal", "Abnormal"],
                            "Count": [int(normal.sum()), int(abnormal.sum())],
                        }
                    )

                fig = px.pie(
                    df_plot,
                    values="Count",
                    names="Category",
                    hole=0.55,
                    title=label,
                )
                fig.update_traces(textinfo="percent+label", textfont_size=18)
                fig.update_layout(showlegend=False, margin=dict(l=0, r=0, t=40, b=0))
                return fig

            def build_population_ppt(
                view_df: pd.DataFrame,
                df_full: pd.DataFrame,
                cols_map: dict,
                location_title: str,
                year_title: str,
            ) -> io.BytesIO:
                """
                Combined A + B style PPT:
                - Keeps all text / bullet analytics
                - Adds image slides for all key Plotly charts across:
                    Overview, Clinical, Lifestyle, Community, Socioeconomic & Profession
                """
                prs = Presentation()

                # Slide Layout Indices (Standard Template)
                # 0: Title, 1: Title+Content, 5: Title Only, 6: Blank
                SLIDE_TITLE = 0
                SLIDE_BULLET = 1
                SLIDE_TITLE_ONLY = 5
                SLIDE_BLANK = 6
                
                # Initialize styling
                plotly_layout = configure_plotly_theme()

                def add_slide(layout_idx):
                    return prs.slides.add_slide(prs.slide_layouts[layout_idx])
                
                def format_title(slide, text, subtitle=None, is_dark=False):
                    """Standardizes title formatting."""
                    title = slide.shapes.title
                    title.text = text
                    tf = title.text_frame
                    p = tf.paragraphs[0]
                    p.font.name = DesignTheme.HEAD_FONT
                    p.font.size = Pt(36)
                    p.font.bold = True
                    
                    if is_dark:
                        p.font.color.rgb = DesignTheme.TEXT_LIGHT
                        # Add a colored background rectangle if it's a "Dark" slide concept
                        bg = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE, 
                            Inches(0), Inches(0), Inches(10), Inches(7.5)
                        )
                        bg.fill.solid()
                        bg.fill.fore_color.rgb = DesignTheme.DARK_BG
                        bg.line.fill.background() # No border
                        pass
                    else:
                        p.font.color.rgb = DesignTheme.PRIMARY
                        
                    if subtitle and len(slide.placeholders) > 1:
                        sub = slide.placeholders[1]
                        sub.text = subtitle
                        sp = sub.text_frame.paragraphs[0]
                        sp.font.name = DesignTheme.BODY_FONT
                        sp.font.size = Pt(20)
                        sp.font.color.rgb = DesignTheme.SECONDARY

                def add_styled_bullets(slide, items, title=None):
                    """Adds specific bullet points with custom styling."""
                    if title:
                        format_title(slide, title)
                        
                    # If standard layout
                    if len(slide.placeholders) > 1:
                        body = slide.placeholders[1].text_frame
                        body.clear()
                        
                        for item in items:
                            p = body.add_paragraph()
                            p.text = item
                            p.font.name = DesignTheme.BODY_FONT
                            p.font.size = Pt(20)
                            p.font.color.rgb = DesignTheme.TEXT_MAIN
                            p.space_after = Pt(10)

                def add_kpi_cards(slide, kpis):
                    """
                    Draws rectangular cards for Key Performance Indicators.
                    kpis: list of dict {'label': str, 'value': str, 'color': RGBColor}
                    """
                    # Start content below title
                    start_y = Inches(2.0)
                    card_width = Inches(2)
                    card_height = Inches(1.5)
                    gap = Inches(0.5)
                    
                    # Center the group of cards
                    total_width = len(kpis) * card_width + (len(kpis) - 1) * gap
                    start_x = (Inches(10) - total_width) / 2
                    
                    for i, kpi in enumerate(kpis):
                        x = start_x + i * (card_width + gap)
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, card_height
                        )
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = kpi.get('color', DesignTheme.SECONDARY)
                        shape.line.color.rgb = DesignTheme.PRIMARY
                        
                        tf = shape.text_frame
                        tf.clear()
                        
                        p_val = tf.add_paragraph()
                        p_val.text = kpi['value']
                        p_val.alignment = PP_ALIGN.CENTER
                        p_val.font.bold = True
                        p_val.font.size = Pt(28)
                        p_val.font.color.rgb = DesignTheme.LIGHT_BG
                        
                        p_lbl = tf.add_paragraph()
                        p_lbl.text = kpi['label']
                        p_lbl.alignment = PP_ALIGN.CENTER
                        p_lbl.font.size = Pt(14)
                        p_lbl.font.color.rgb = DesignTheme.LIGHT_BG

                def add_plotly_image(slide, fig, title=None):
                    """Adds a Plotly figure as an image."""
                    if title:
                        format_title(slide, title)
                        
                    fig.update_layout(plotly_layout)
                    fig.update_traces(textfont_size=18)
                    fig.update_layout(margin=dict(l=20, r=20, t=50, b=20))
                    
                    img_bytes = fig.to_image(format="png", width=1200, height=750, scale=2)
                    
                    # Center image
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9.0)
                    height = Inches(5.5)
                    
                    slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)

                def add_donut_chart(slide, labels, values, left, top, width, height, title=None):
                    fig = px.pie(
                        names=labels,
                        values=values,
                        hole=0.55,
                        color_discrete_sequence=px.colors.qualitative.Set2
                    )

                    fig.update_traces(
                        textposition="outside",
                        textinfo="percent+label",
                        textfont=dict(
                            size=34,
                            color="#444444"
                        ),
                        marker=dict(line=dict(color="white", width=2))
                    )

                    fig.update_layout(
                        showlegend=False,
                        paper_bgcolor="white",
                        plot_bgcolor="white",
                        margin=dict(l=30, r=30, t=40, b=30),
                        font=dict(size=26, color="#444444"),
                    )

                    if title:
                        fig.update_layout(
                            title=dict(text=title, font=dict(size=28, color="#333333"), x=0.5)
                        )

                    # Export a much larger, sharper image
                    img = fig.to_image(format="png", width=1200, height=900, scale=2)

                    slide.shapes.add_picture(io.BytesIO(img), left, top, width=width, height=height)

                # ---------------- Common metrics -----------------

                total_screened = len(view_df)
                total_dataset = len(df_full)

                if "Health Status" in view_df.columns:
                    normal_mask = view_df["Health Status"] == "Healthy"
                    total_normal = int(normal_mask.sum())
                    total_abnormal = int(total_screened - total_normal)
                else:
                    total_normal = 0
                    total_abnormal = 0

                pct_normal = safe_pct(total_normal, total_screened)
                pct_abnormal = safe_pct(total_abnormal, total_screened)

                cols = cols_map

                # ---------------- Slide 1: Title -----------------

                slide = add_slide(SLIDE_TITLE)
                
                # Custom Background for Title
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
                bg.fill.solid()
                bg.fill.fore_color.rgb = DesignTheme.DARK_BG
                bg.line.fill.background()
                
                # Manually add text boxes on top to ensure visibility
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
                tf = title_box.text_frame
                p = tf.add_paragraph()
                p.text = f"Population Health Assessment\n{location_title}"
                p.alignment = PP_ALIGN.CENTER
                p.font.name = DesignTheme.HEAD_FONT
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = DesignTheme.LIGHT_BG
                
                sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
                tf = sub_box.text_frame
                p = tf.add_paragraph()
                p.text = f"Screening Summary • {year_title}\nGenerated via iHealth_Map"
                p.alignment = PP_ALIGN.CENTER
                p.font.name = DesignTheme.BODY_FONT
                p.font.size = Pt(20)
                p.font.color.rgb = DesignTheme.SECONDARY

                # ---------------- Slide 2: Overall Screening Summary -----------------

                slide = add_slide(SLIDE_TITLE_ONLY)
                format_title(slide, "Overall Screening Summary")

                kpis = [
                    {"label": "Dataset Size", "value": f"{total_dataset:,}", "color": DesignTheme.PRIMARY},
                    {"label": "Screened", "value": f"{total_screened:,}", "color": DesignTheme.SECONDARY},
                    {"label": "Abnormal", "value": f"{pct_abnormal:.1f}%", "color": DesignTheme.ACCENT},
                    {"label": "Normal", "value": f"{pct_normal:.1f}%", "color": DesignTheme.SECONDARY},
                ]
                add_kpi_cards(slide, kpis)

                summary_lines = [
                    f"Total persons in dataset: {total_dataset:,}",
                    f"Screened persons (after filters): {total_screened:,}",
                    f"Total with abnormal health status: {total_abnormal:,} ({pct_abnormal:.1f}%)",
                    f"Total with normal health status: {total_normal:,} ({pct_normal:.1f}%)",
                ]
                add_styled_bullets(slide, summary_lines)

                # ---------------- Slide 3: Gender-wise Participation -----------------

                if cols["gender"]:
                    g_series = view_df[cols["gender"]].astype(str)
                    g_counts = g_series.value_counts()
                    total_g = int(g_counts.sum())
                    lines = [
                        f"{g}: {cnt:,} ({safe_pct(cnt, total_g):.1f}%)"
                        for g, cnt in g_counts.items()
                    ]
                else:
                    lines = ["Gender column not available in this dataset."]

                slide = add_slide(SLIDE_BULLET)
                add_styled_bullets(slide, lines, title="Gender-wise Participation")

                # --- Add donut chart below bullets ---
                if cols["gender"]:
                    labels = g_counts.index.tolist()
                    values = g_counts.values.tolist()

                    add_donut_chart(
                        slide,
                        labels=labels,
                        values=values,
                        left=Inches(2.5),
                        top=Inches(3.5),
                        width=Inches(5),
                        height=Inches(3.6),
                        title=None
                    )

                # -----------------------------------------------------------------
                # Parameter abnormality helpers (logic unchanged)
                # -----------------------------------------------------------------

                glu_col = cols["glucose"]
                chol_col = cols["chol"]
                bp_sys = cols["bp_sys"]
                bp_dia = cols["bp_dia"]
                creat_col = cols["creatinine"]
                alt_col = cols["alt"]
                prot_col = cols["protein_total"]
                bmi_col = cols["bmi"]
                locality_col = cols["locality"]

                def abnormal_bp_mask(df_in):
                    if not (bp_sys and bp_dia):
                        return pd.Series([False] * len(df_in))
                    s = as_num(df_in[bp_sys])
                    d = as_num(df_in[bp_dia])
                    return (s >= 130) | (d >= 80)

                def abnormal_glu_mask(df_in):
                    if not glu_col:
                        return pd.Series([False] * len(df_in))
                    g = as_num(df_in[glu_col])
                    return g >= 126

                def abnormal_chol_mask(df_in):
                    if not chol_col:
                        return pd.Series([False] * len(df_in))
                    ch = as_num(df_in[chol_col])
                    return ch >= 240

                def abnormal_creat_mask(df_in):
                    if not creat_col:
                        return pd.Series([False] * len(df_in))
                    c = as_num(df_in[creat_col])
                    return c > 1.2

                def abnormal_protein_mask(df_in):
                    if not prot_col:
                        return pd.Series([False] * len(df_in))
                    p = as_num(df_in[prot_col])
                    return (p < 6.0) | (p > 8.0)

                def abnormal_alt_mask(df_in):
                    if not alt_col:
                        return pd.Series([False] * len(df_in))
                    a = as_num(df_in[alt_col])
                    return a > 40

                # ---------------- Slides: Abnormality summaries ----------------

                def abnormal_summary(title_text, series, predicate):
                    valid = series.dropna()
                    if len(valid) == 0:
                        slide = add_slide(SLIDE_BULLET)
                        add_styled_bullets(slide, ["No data available."], title=title_text)
                        return

                    ab = predicate(valid)
                    cnt_ab = int(ab.sum())
                    cnt_norm = len(valid) - cnt_ab

                    slide = add_slide(SLIDE_BULLET)
                    add_styled_bullets(
                        slide,
                        [
                            f"Abnormal: {cnt_ab:,} ({safe_pct(cnt_ab, len(valid)):.1f}%)",
                            f"Normal: {cnt_norm:,} ({safe_pct(cnt_norm, len(valid)):.1f}%)",
                        ],
                        title=title_text,
                    )

                    # --- Donut chart below bullets ---
                    add_donut_chart(
                        slide,
                        labels=["Abnormal", "Normal"],
                        values=[cnt_ab, cnt_norm],
                        left=Inches(2.5),
                        top=Inches(3.5),
                        width=Inches(5),
                        height=Inches(3.6),
                        title=None
                    )

                if glu_col:
                    abnormal_summary(
                        "Diabetes (fasting glucose ≥126 mg/dL)",
                        as_num(view_df[glu_col]),
                        lambda x: x >= 126,
                    )

                if bp_sys and bp_dia:
                    valid_bp = view_df[[bp_sys, bp_dia]].dropna()
                    if len(valid_bp) == 0:
                        slide = add_slide(SLIDE_BULLET)
                        add_styled_bullets(slide, ["No BP data available."], title="Blood Pressure Status")
                    else:
                        ab = abnormal_bp_mask(valid_bp)
                        cnt_ab = int(ab.sum())
                        cnt_norm = len(valid_bp) - cnt_ab

                        slide = add_slide(SLIDE_BULLET)
                        add_styled_bullets(
                            slide,
                            [
                                f"Hypertensive (>130/80): {cnt_ab:,} ({safe_pct(cnt_ab, len(valid_bp)):.1f}%)",
                                f"Normotensive: {cnt_norm:,} ({safe_pct(cnt_norm, len(valid_bp)):.1f}%)",
                            ],
                            title="Blood Pressure Status",
                        )

                        # --- Donut chart below bullets ---
                        add_donut_chart(
                            slide,
                            labels=["Hypertensive", "Normotensive"],
                            values=[cnt_ab, cnt_norm],
                            left=Inches(2.5),
                            top=Inches(3.5),
                            width=Inches(5),
                            height=Inches(3.6),
                            title=None
                        )

                if chol_col:
                    abnormal_summary(
                        "Cholesterol (≥240 mg/dL)",
                        as_num(view_df[chol_col]),
                        lambda x: x >= 240,
                    )

                if creat_col:
                    abnormal_summary(
                        "Renal (Creatinine >1.2 mg/dL)",
                        as_num(view_df[creat_col]),
                        lambda x: x > 1.2,
                    )

                if prot_col:
                    abnormal_summary(
                        "Total Protein (outside 6.0–8.0 g/dL)",
                        as_num(view_df[prot_col]),
                        lambda x: (x < 6.0) | (x > 8.0),
                    )

                if alt_col:
                    abnormal_summary(
                        "Liver Enzyme (ALT/SGPT >40 U/L)",
                        as_num(view_df[alt_col]),
                        lambda x: x > 40,
                    )

                # ---------------- Age 18–20 yrs abnormality -------------------

                if "__AGE__" in view_df.columns:
                    grp_1820 = view_df[view_df["__AGE__"].between(18, 20)]

                    slide = add_slide(SLIDE_BULLET)
                    if len(grp_1820) == 0:
                        add_styled_bullets(slide, ["No records in the 18–20 years age group."],
                                        title="Abnormal Parameters – 18–20 years")
                    else:
                        bullets = []
                        for label, fn in [
                            ("Hypertension", abnormal_bp_mask),
                            ("Diabetes (≥126 mg/dL)", abnormal_glu_mask),
                            ("High cholesterol (≥240 mg/dL)", abnormal_chol_mask),
                            ("Abnormal creatinine", abnormal_creat_mask),
                            ("Abnormal total protein", abnormal_protein_mask),
                            ("Abnormal SGPT", abnormal_alt_mask),
                        ]:
                            ab = fn(grp_1820)
                            bullets.append(f"{label}: {safe_pct(ab.sum(), len(grp_1820)):.1f}%")

                        add_styled_bullets(slide, bullets, title="Abnormal Parameters – 18–20 years")
                # else:
                #     slide = add_slide(SLIDE_BULLET)
                #     add_styled_bullets(slide, ["Age column not available."],
                #                     title="Abnormal Parameters – 18–20 years")

                # ---------------- Community-wise abnormal parameters ----------

                if locality_col:
                    slide = add_slide(SLIDE_BULLET)
                    sub_loc = view_df.copy()
                    any_ab = (
                        abnormal_bp_mask(sub_loc)
                        | abnormal_glu_mask(sub_loc)
                        | abnormal_chol_mask(sub_loc)
                        | abnormal_creat_mask(sub_loc)
                        | abnormal_protein_mask(sub_loc)
                        | abnormal_alt_mask(sub_loc)
                    )

                    sub_loc["__ANY_ABNORMAL__"] = any_ab.astype(int)
                    grp = (
                        sub_loc.groupby(locality_col)["__ANY_ABNORMAL__"]
                        .agg(["count", "sum"])
                        .reset_index()
                    )
                    grp["abnormal_pct"] = grp["sum"] * 100.0 / grp["count"]
                    grp = grp.sort_values("count", ascending=False).head(10)

                    bullets = []
                    for _, row in grp.iterrows():
                        name = str(row[locality_col])
                        n_pop = int(row["count"])
                        p_ab = float(row["abnormal_pct"])
                        bullets.append(f"{name}: {p_ab:.1f}% with ≥1 abnormal parameter (n={n_pop})")

                    add_styled_bullets(slide, bullets, title="Abnormal Parameters – Community-wise")
                # else:
                #     add_styled_bullets(slide, ["Community/locality column not available."],
                #                     title="Abnormal Parameters – Community-wise")

                # --- Histogram slide: Community-wise abnormality ---
                if locality_col:

                    slide = add_slide(SLIDE_BULLET)
                    sub_loc = view_df.copy()
                    any_ab = (
                        abnormal_bp_mask(sub_loc)
                        | abnormal_glu_mask(sub_loc)
                        | abnormal_chol_mask(sub_loc)
                        | abnormal_creat_mask(sub_loc)
                        | abnormal_protein_mask(sub_loc)
                        | abnormal_alt_mask(sub_loc)
                    )

                    sub_loc["__ANY_ABNORMAL__"] = any_ab.astype(int)
                    grp = (
                        sub_loc.groupby(locality_col)["__ANY_ABNORMAL__"]
                        .agg(["count", "sum"])
                        .reset_index()
                    )
                    grp["abnormal_pct"] = grp["sum"] * 100.0 / grp["count"]
                    grp = grp.sort_values("count", ascending=False).head(10)

                    bullets = []
                    for _, row in grp.iterrows():
                        name = str(row[locality_col])
                        n_pop = int(row["count"])
                        p_ab = float(row["abnormal_pct"])
                        bullets.append(f"{name}: {p_ab:.1f}% with ≥1 abnormal parameter (n={n_pop})")

                    add_styled_bullets(slide, bullets, title="Abnormal Parameters – Community-wise")

                    # --- Histogram slide: Community-wise abnormality ---
                    hist_slide = add_slide(SLIDE_TITLE_ONLY)
                    format_title(hist_slide, "Abnormal Parameters – Community-wise")

                    fig = px.bar(
                        grp,
                        x=locality_col,
                        y="abnormal_pct",
                        text=grp["abnormal_pct"].round(1),
                        labels={
                            locality_col: "Community",
                            "abnormal_pct": "% Population with ≥1 Abnormal Parameter",
                        },
                    )

                    fig.update_traces(textposition="outside")
                    fig.update_layout(
                        yaxis_range=[0, max(100, grp["abnormal_pct"].max() + 10)],
                        height=600,
                        margin=dict(t=60, b=60, l=60, r=60),
                        paper_bgcolor="white",
                        plot_bgcolor="white",
                        font=dict(color="#444444", size=18),
                        title_font=dict(color="#444444", size=22),
                        xaxis=dict(showgrid=False, linecolor="#444444",
                                tickfont=dict(color="#444444"),
                                title_font=dict(color="#444444")),
                        yaxis=dict(showgrid=True, gridcolor="#dddddd",
                                linecolor="#444444",
                                tickfont=dict(color="#444444"),
                                title_font=dict(color="#444444")),
                    )

                    img_bytes = fig.to_image(format="png", scale=2)

                    width = Inches(6.4)
                    left = Inches(1.8)
                    top = Inches(1.6)

                    hist_slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width)

                # ---------------- BMI distribution --------------------

                if bmi_col:
                    slide = add_slide(SLIDE_BULLET)
                    b_ = as_num(view_df[bmi_col]).dropna()
                    n_b = len(b_)

                    if n_b == 0:
                        add_styled_bullets(slide, ["No BMI data available."], title="BMI Distribution")
                    else:
                        under = (b_ < 18.5).sum()
                        normal_b = ((b_ >= 18.5) & (b_ < 25)).sum()
                        over = ((b_ >= 25) & (b_ < 30)).sum()
                        obese = (b_ >= 30).sum()

                        lines = [
                            f"Underweight (<18.5): {under:,} ({safe_pct(under, n_b):.1f}%)",
                            f"Normal (18.5–24.9): {normal_b:,} ({safe_pct(normal_b, n_b):.1f}%)",
                            f"Overweight (25–29.9): {over:,} ({safe_pct(over, n_b):.1f}%)",
                            f"Obese (≥30): {obese:,} ({safe_pct(obese, n_b):.1f}%)",
                            "Classification: <18.5 underweight; 18.5–24.9 healthy; 25–29.9 overweight; ≥30 obese.",
                        ]

                        add_styled_bullets(slide, lines, title="BMI Distribution")
                # else:
                #     add_styled_bullets(slide, ["BMI column not available."], title="BMI Distribution")

                # ---------------- BMI vs age bands --------------------

                if bmi_col and "__AGE__" in view_df.columns:
                    slide = add_slide(SLIDE_BULLET)
                    sub = view_df[["__AGE__", bmi_col]].dropna()

                    if len(sub) > 0:
                        bins = [0, 10, 20, 30, 40, 50, 60, 70, 200]
                        labels = ["01–10","11–20","21–30","31–40","41–50","51–60","61–70","71+"]

                        sub["age_band"] = pd.cut(sub["__AGE__"], bins=bins, labels=labels, right=True)
                        b_ = as_num(sub[bmi_col])
                        sub["over_obese"] = (b_ >= 25).astype(int)

                        grp = (
                            sub.groupby("age_band")["over_obese"]
                            .agg(["count", "sum"])
                            .reset_index()
                        )
                        grp["pct_over_obese"] = grp["sum"] * 100.0 / grp["count"]
                        grp = grp.dropna(subset=["age_band"])

                        lines = []
                        for _, row in grp.iterrows():
                            lines.append(f"{row['age_band']} years: {row['pct_over_obese']:.1f}% overweight/obese")

                        peak = grp.loc[grp["pct_over_obese"].idxmax(), "age_band"]
                        lines.append(f"Peak overweight/obesity seen in {peak} years group.")

                        add_styled_bullets(slide, lines, title="BMI in Different Age Groups")
                    else:
                        add_styled_bullets(slide, ["No BMI+age records available."],
                                        title="BMI in Different Age Groups")
                # else:
                #     add_styled_bullets(slide, ["BMI or age column not available."],
                #                     title="BMI in Different Age Groups")

                # ---------------- Rural vs Urban Comparison ------------

                if locality_col:
                    slide = add_slide(SLIDE_BULLET)
                    loc_series = view_df[locality_col].astype(str).str.lower()

                    def rural_urban_tag(val):
                        if "indore" in val: return "Urban (Indore)"
                        if "mhow" in val: return "Rural (Mhow)"
                        return "Other"

                    df_ru = view_df.copy()
                    df_ru["__RU_TAG__"] = loc_series.map(rural_urban_tag)

                    metrics = []
                    for tag in ["Rural (Mhow)", "Urban (Indore)"]:
                        sub = df_ru[df_ru["__RU_TAG__"] == tag]
                        if len(sub):
                            metrics.append(
                                f"{tag}: n={len(sub)}; "
                                f"mean fasting glucose ≈ {as_num(sub[glu_col]).mean():.1f} mg/dL; "
                                f"mean cholesterol ≈ {as_num(sub[chol_col]).mean():.1f} mg/dL; "
                                f"mean SGPT ≈ {as_num(sub[alt_col]).mean():.1f} U/L"
                            )

                    if metrics:
                        add_styled_bullets(slide, metrics, title="Rural vs Urban Comparison")
                    # else:
                    #     add_styled_bullets(slide, [
                    #         "Rural (Mhow) vs Urban (Indore) tags could not be derived.",
                    #         "Ensure locality text contains 'Mhow' and 'Indore'.",
                    #     ], title="Rural vs Urban Comparison")
                # else:
                #     add_styled_bullets(slide, ["Locality column not available; comparison not generated."],
                #                     title="Rural vs Urban Comparison")

                # ---------------- Save PPT -------------------------
                buf = io.BytesIO()
                prs.save(buf)
                buf.seek(0)
                return buf


            # Button to generate PPT + show download link
            if st.button("Generate PPT summary"):
                st.session_state.ppt_version += 1
                st.session_state.ppt_buffer = build_population_ppt(
                    view,
                    df,
                    cols,
                    ppt_location,
                    ppt_year,
                )

            if st.session_state.ppt_buffer is not None:
                st.download_button(
                    "⬇️ Download PPT summary",
                    data=st.session_state.ppt_buffer,
                    file_name="bharat_ihealthmap_summary.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

        # -----------------------------------------------------------------
        #  Individual PDF reports (unchanged)
        # -----------------------------------------------------------------
        st.markdown("---")
        st.subheader("Generate individual patient report (PDF)")

        id_label = ID_COL
        st.caption(f"Primary ID column: **{id_label}**")

        id_input = st.text_input(
            "Enter one or more IDs (comma-separated) from the primary ID column",
            help="Example: LAB00123, LAB00456",
        )

        def build_patient_pdf(df_base: pd.DataFrame, ids: List[str]) -> io.BytesIO:
            # --- ID normalisation for reliable matching ---
            df_match = df_base.copy()
            df_match["_clean_id"] = (
                df_match[id_label]
                .astype(str)
                .str.strip()
                .str.replace(r"\.0$", "", regex=True)
                .str.replace(r"\s+", "", regex=True)
                .str.upper()
            )

            cleaned_ids = [
                re.sub(r"\s+", "", pid.strip()).upper().replace(".0", "")
                for pid in ids
            ]

            buf = io.BytesIO()
            c = canvas.Canvas(buf, pagesize=letter)
            width, height = letter
            c.setFont("Helvetica", 10)

            for cid in cleaned_ids:
                sub = df_match[df_match["_clean_id"] == cid]

                if sub.empty:
                    c.setFont("Helvetica-Bold", 14)
                    c.drawString(40, height - 40, f"ID: {cid} not found in dataset.")
                    c.showPage()
                    continue

                row = sub.iloc[0]

                # Header
                y = height - 40
                c.setFont("Helvetica-Bold", 14)
                c.drawString(40, y, "CentralLab iHealthMAP – Individual Health Report")
                y -= 24

                c.setFont("Helvetica", 10)
                c.drawString(40, y, f"ID: {cid}")
                y -= 14

                # Basic details
                if cols["name"]:
                    c.drawString(40, y, f"Name: {row.get(cols['name'], '')}")
                    y -= 14
                if "__AGE__" in row:
                    c.drawString(40, y, f"Age: {row.get('__AGE__', '')}")
                    y -= 14
                if cols["gender"]:
                    c.drawString(40, y, f"Gender: {row.get(cols['gender'], '')}")
                    y -= 14
                if cols["locality"]:
                    c.drawString(40, y, f"Locality: {row.get(cols['locality'], '')}")
                    y -= 14

                # Measurements
                y -= 10
                c.setFont("Helvetica-Bold", 11)
                c.drawString(40, y, "Key Measurements")
                y -= 16
                c.setFont("Helvetica", 10)

                def put_line(label, value):
                    nonlocal y
                    if y < 60:
                        c.showPage()
                        y = height - 40
                        c.setFont("Helvetica", 10)
                    c.drawString(50, y, f"{label}: {value}")
                    y -= 14

                if cols["bmi"]:
                    put_line("BMI", row.get(cols["bmi"], ""))
                if cols["glucose"]:
                    put_line("Fasting glucose (mg/dL)", row.get(cols["glucose"], ""))
                if cols["chol"]:
                    put_line("Total cholesterol (mg/dL)", row.get(cols["chol"], ""))
                if cols["bp_sys"] and cols["bp_dia"]:
                    put_line(
                        "Blood pressure",
                        f"{row.get(cols['bp_sys'], '')}/{row.get(cols['bp_dia'], '')}",
                    )
                if cols["creatinine"]:
                    put_line("Creatinine (mg/dL)", row.get(cols["creatinine"], ""))
                if cols["alt"]:
                    put_line("ALT (SGPT)", row.get(cols["alt"], ""))

                # Assessment Block
                y -= 10
                c.setFont("Helvetica-Bold", 11)
                c.drawString(40, y, "Automated Assessment")
                y -= 18
                c.setFont("Helvetica", 10)

                def put_block(label, text):
                    nonlocal y
                    put_line(label, "")
                    words = str(text).split()
                    line = ""
                    for w in words:
                        if len(line) + len(w) + 1 > 95:
                            put_line("   " + line, "")
                            line = w
                        else:
                            line += (" " if line else "") + w
                    if line:
                        put_line("   " + line, "")

                put_block("Health Status", row.get("Health Status", ""))
                put_block("Diagnosis", row.get("Diagnosis", ""))
                put_block("Prognosis", row.get("Prognosis", ""))
                put_block("Recommendations", row.get("Prescription", ""))

                c.showPage()

            c.save()
            buf.seek(0)
            return buf

        if st.button("Generate PDF report(s)") and id_input.strip():
            ids = [s.strip() for s in id_input.split(",") if s.strip()]
            if ids:
                pdf_buf = build_patient_pdf(view, ids)
                st.download_button(
                    "⬇️ Download generated PDF",
                    data=pdf_buf,
                    file_name="central_lab_ihealthmap_reports.pdf",
                    mime="application/pdf",
                )

# ---------------------------------------------------------------------
#  Footer
# ---------------------------------------------------------------------
st.markdown("---")
st.caption(
    "CentralLab iHealthMAP – Prototype population health analytic dashboard. "
    "Screening thresholds approximate ADA/WHO-style cutoffs and do not replace clinical diagnosis."
)
