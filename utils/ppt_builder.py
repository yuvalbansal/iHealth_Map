import io
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import plotly.express as px
import plotly.graph_objects as go

# ---------------------------------------------------------------------
# Design & Theme Configuration
# ---------------------------------------------------------------------

class DesignTheme:
    """
    Centralized theme configuration for the presentation.
    Theme: Clean Medical/Tech
    """
    # Colors (RGB)
    PRIMARY = RGBColor(0, 91, 150)       # Deep Teal/Blue #005b96
    SECONDARY = RGBColor(100, 151, 177)  # Soft Blue #6497b1
    ACCENT = RGBColor(255, 111, 105)     # Coral/Orange #ff6f69
    DARK_BG = RGBColor(3, 37, 76)        # Dark Night Blue #03254c (for Titles)
    LIGHT_BG = RGBColor(255, 255, 255)   # White
    TEXT_MAIN = RGBColor(50, 50, 50)     # Dark Grey
    TEXT_LIGHT = RGBColor(240, 240, 240) # Off-white for dark backgrounds
    
    # Fonts
    HEAD_FONT = "Arial"
    BODY_FONT = "Arial"
    
    # Plotly Template Name
    PLOTLY_TEMPLATE = "plotly_white"
    
    # Color Sequence for Charts
    COLOR_SEQUENCE = ["#005b96", "#6497b1", "#ff6f69", "#03254c", "#b3cde0"]

def configure_plotly_theme():
    """Returns a layout template for Plotly charts to match PPT style."""
    return dict(
        font_family=DesignTheme.BODY_FONT,
        font_color="#333333",
        font_size=16,
        title_font_family=DesignTheme.HEAD_FONT,
        title_font_size=24,
        colorway=DesignTheme.COLOR_SEQUENCE,
        plot_bgcolor="rgba(0,0,0,0)",
        paper_bgcolor="rgba(0,0,0,0)",
        xaxis=dict(
            showgrid=False, 
            linecolor="#333",
            title_font_size=20,
            tickfont_size=16
        ),
        yaxis=dict(
            showgrid=True, 
            gridcolor="#eee", 
            linecolor="#333",
            title_font_size=20,
            tickfont_size=16
        ),
        legend=dict(
            font_size=18
        ),
    )

# ---------------------------------------------------------------------
# Helper Utilities
# ---------------------------------------------------------------------

def as_num(series: pd.Series) -> pd.Series:
    return pd.to_numeric(series, errors="coerce")


def safe_pct(num: float, den: float) -> float:
    return float(num) * 100.0 / float(den) if den not in (0, None) else 0.0


def set_autofit(tf, maxsize=28, minsize=12):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    try:
        tf.fit_text(max_size=maxsize, min_size=minsize)
    except Exception:
        pass


def donut_normal_abnormal(series, label, thresholds):
    s = pd.to_numeric(series, errors="coerce")
    s = s[s.notna()]

    if len(s) == 0:
        df_plot = pd.DataFrame({"Category": ["No data"], "Count": [1]})
    else:
        low_ok, high_ok = thresholds
        if low_ok is not None and high_ok is not None:
            normal = s.between(low_ok, high_ok)
        elif low_ok is None and high_ok is not None:
            normal = s <= high_ok
        elif low_ok is not None and high_ok is None:
            normal = s >= low_ok
        else:
            normal = pd.Series(False, index=s.index)

        abnormal = ~normal
        df_plot = pd.DataFrame(
            {
                "Category": ["Normal", "Abnormal"],
                "Count": [int(normal.sum()), int(abnormal.sum())],
            }
        )

    fig = px.pie(
        df_plot,
        values="Count",
        names="Category",
        hole=0.55,
        title=label,
    )
    fig.update_traces(textinfo="percent+label", textfont_size=18)
    fig.update_layout(showlegend=False, margin=dict(l=0, r=0, t=40, b=0))
    return fig


# ---------------------------------------------------------------------
# Main public API
# ---------------------------------------------------------------------

def build_population_ppt(
    view_df: pd.DataFrame,
    df_full: pd.DataFrame,
    cols_map: dict,
    location_title: str,
    year_title: str,
) -> io.BytesIO:
    """
    Generates a beautifully formatted Population Health PPT.
    """
    prs = Presentation()
    
    # Slide Layout Indices (Standard Template)
    # 0: Title, 1: Title+Content, 5: Title Only, 6: Blank
    SLIDE_TITLE = 0
    SLIDE_BULLET = 1
    SLIDE_TITLE_ONLY = 5
    SLIDE_BLANK = 6
    
    # Initialize styling
    plotly_layout = configure_plotly_theme()
    
    # --- Internal Helper Functions for this PPT ---
    
    def add_slide(layout_idx):
        return prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    def format_title(slide, text, subtitle=None, is_dark=False):
        """Standardizes title formatting."""
        title = slide.shapes.title
        title.text = text
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.font.name = DesignTheme.HEAD_FONT
        p.font.size = Pt(36)
        p.font.bold = True
        
        if is_dark:
            p.font.color.rgb = DesignTheme.TEXT_LIGHT
             # Add a colored background rectangle if it's a "Dark" slide concept
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(0), Inches(0), Inches(10), Inches(7.5)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = DesignTheme.DARK_BG
            bg.line.fill.background() # No border
            # Send to back hack: usually new shapes are on top. 
            # We must create bg FIRST if we want it behind, OR move title to front.
            # Easier approach: Create BG slide first, then add text boxes.
            # BUT: In python-pptx, 'title' is a placeholder. 
            # We'll just set the title color and assume the user uses a template 
            # or we accept that 'Dark' means we manually add a background shape BEHIND.
            # Since z-ordering is tricky, let's keep it simple: 
            # standard white slides with colored headers.
            pass
        else:
            p.font.color.rgb = DesignTheme.PRIMARY
            
        if subtitle and len(slide.placeholders) > 1:
            sub = slide.placeholders[1]
            sub.text = subtitle
            sp = sub.text_frame.paragraphs[0]
            sp.font.name = DesignTheme.BODY_FONT
            sp.font.size = Pt(20)
            sp.font.color.rgb = DesignTheme.SECONDARY

    def add_styled_bullets(slide, items, title=None):
        """Adds specific bullet points with custom styling."""
        if title:
            format_title(slide, title)
            
        # If standard layout
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()
            
            for item in items:
                p = body.add_paragraph()
                p.text = item
                p.font.name = DesignTheme.BODY_FONT
                p.font.size = Pt(20)
                p.font.color.rgb = DesignTheme.TEXT_MAIN
                p.space_after = Pt(10)

    def add_kpi_cards(slide, kpis):
        """
        Draws rectangular cards for Key Performance Indicators.
        kpis: list of dict {'label': str, 'value': str, 'color': RGBColor}
        """
        # Start content below title
        start_y = Inches(2.0)
        card_width = Inches(2.5)
        card_height = Inches(1.5)
        gap = Inches(0.5)
        
        # Center the group of cards
        total_width = len(kpis) * card_width + (len(kpis) - 1) * gap
        start_x = (Inches(10) - total_width) / 2
        
        for i, kpi in enumerate(kpis):
            x = start_x + i * (card_width + gap)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, card_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = kpi.get('color', DesignTheme.SECONDARY)
            shape.line.color.rgb = DesignTheme.PRIMARY
            
            tf = shape.text_frame
            tf.clear()
            
            p_val = tf.add_paragraph()
            p_val.text = kpi['value']
            p_val.alignment = PP_ALIGN.CENTER
            p_val.font.bold = True
            p_val.font.size = Pt(28)
            p_val.font.color.rgb = DesignTheme.LIGHT_BG
            
            p_lbl = tf.add_paragraph()
            p_lbl.text = kpi['label']
            p_lbl.alignment = PP_ALIGN.CENTER
            p_lbl.font.size = Pt(14)
            p_lbl.font.color.rgb = DesignTheme.LIGHT_BG

    def add_plotly_image(slide, fig, title=None):
        """Adds a Plotly figure as an image."""
        if title:
            format_title(slide, title)
            
        fig.update_layout(plotly_layout)
        fig.update_traces(textfont_size=18)
        fig.update_layout(margin=dict(l=20, r=20, t=50, b=20))
        
        img_bytes = fig.to_image(format="png", width=1200, height=750, scale=2)
        
        # Center image
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(5.5)
        
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)

    # --- Data Prep Helpers ---
    def pct(n, d): return round((n / d) * 100, 1) if d else 0.0
    
    # ---------------- Slide 1: Title Slide ----------------
    slide = add_slide(SLIDE_TITLE)
    
    # Custom Background for Title
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DesignTheme.DARK_BG
    bg.line.fill.background()
    
    # Manually add text boxes on top to ensure visibility
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = f"Population Health Assessment\n{location_title}"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = DesignTheme.HEAD_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DesignTheme.LIGHT_BG
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = f"Screening Summary • {year_title}\nGenerated via iHealth_Map"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = DesignTheme.BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = DesignTheme.SECONDARY

    # ---------------- Slide 2: Dataset & Coverage (KPIs) ----------------
    slide = add_slide(SLIDE_TITLE_ONLY)
    format_title(slide, "Dataset & Coverage Summary")
    
    total_all = len(df_full)
    total_view = len(view_df)
    
    kpis = [
        {"label": "Total Records", "value": f"{total_all:,}", "color": DesignTheme.PRIMARY},
        {"label": "Analyzed", "value": f"{total_view:,}", "color": DesignTheme.SECONDARY},
        {"label": "Coverage", "value": f"{pct(total_view, total_all)}%", "color": DesignTheme.ACCENT},
    ]
    add_kpi_cards(slide, kpis)
    
    # Add text summary below
    txt_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf = txt_box.text_frame
    p = tf.add_paragraph()
    p.level = 0
    p.font.size = Pt(16)
    
    extras = []
    if "__AGE__" in view_df:
        mn, mx = int(view_df["__AGE__"].min()), int(view_df["__AGE__"].max())
        extras.append(f"Age range: {mn}–{mx} years")
    
    gender_c = cols_map.get("gender")
    if gender_c and gender_c in view_df:
        counts = view_df[gender_c].value_counts()
        top_g = counts.index[0] if len(counts) > 0 else "N/A"
        extras.append(f"Dominant Gender: {top_g}")
        
    p.text = "\n".join(extras)

    # ---------------- Slide 3: Overall Health Status ----------------
    # Donut Chart Left, Key Insight Right
    if "Health Status" in view_df.columns:
        slide = add_slide(SLIDE_TITLE_ONLY)
        format_title(slide, "Overall Health Status")
        
        status_counts = view_df["Health Status"].replace({"Needs Attention": "At Risk"}).value_counts().reset_index()
        status_counts.columns = ["Status", "Count"]
        
        fig = px.pie(status_counts, values="Count", names="Status", hole=0.6,
                     color_discrete_sequence=[DesignTheme.COLOR_SEQUENCE[0], DesignTheme.COLOR_SEQUENCE[2]])
        fig.update_traces(textinfo="percent+label", textfont_size=20)
        
        # Render Chart Left
        fig.update_layout(plotly_layout)
        img_bytes = fig.to_image(format="png", width=600, height=500)
        slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(2.0), width=Inches(5))
        
        # Text Right
        tb = slide.shapes.add_textbox(Inches(5.5), Inches(3), Inches(3.5), Inches(3))
        tf = tb.text_frame
        
        at_risk = status_counts.loc[status_counts["Status"] != "Healthy", "Count"].sum()
        risk_pct = pct(at_risk, total_view)
        
        p = tf.add_paragraph()
        p.text = f"{risk_pct}%"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = DesignTheme.ACCENT
        
        p2 = tf.add_paragraph()
        p2.text = "of the population has at least"
        p2.font.size = Pt(20)

        p3 = tf.add_paragraph()
        p3.text = "one abnormal health indicator."
        p3.font.size = Pt(20)

    # ---------------- Slide 4: Key Health Risk Burden ----------------
    slide = add_slide(SLIDE_TITLE_ONLY)
    format_title(slide, "Key Health Risk Burden")
    
    risks = []
    if cols_map.get("glucose"):
        g = as_num(view_df[cols_map["glucose"]])
        val = pct((g >= 126).sum(), g.notna().sum())
        risks.append({"label": "Diabetes", "value": f"{val}%", "color": DesignTheme.SECONDARY})
        
    if cols_map.get("bp_sys") and cols_map.get("bp_dia"):
        s = as_num(view_df[cols_map["bp_sys"]])
        d = as_num(view_df[cols_map["bp_dia"]])
        val = pct(((s>=130)|(d>=80)).sum(), s.notna().sum())
        risks.append({"label": "Hypertension", "value": f"{val}%", "color": DesignTheme.SECONDARY})
        
    if cols_map.get("bmi"):
        b = as_num(view_df[cols_map["bmi"]])
        val = pct((b>=30).sum(), b.notna().sum())
        risks.append({"label": "Obesity", "value": f"{val}%", "color": DesignTheme.ACCENT}) # emphasize obesity
        
    if risks:
        add_kpi_cards(slide, risks)
    else:
        tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tb.text_frame.text = "No clinical data available for risk assessment."

    # ---------------- Slide 5: Age-wise Risk Escalation ----------------
    if "__AGE__" in view_df.columns and "Health Status" in view_df.columns:
        slide = add_slide(SLIDE_TITLE_ONLY)
        
        df_a = view_df.copy()
        df_a["is_bad"] = df_a["Health Status"] != "Healthy"
        bins = [0,18,30,40,50,60,70,200]
        labels = ["<18","18-29","30-39","40-49","50-59","60-69","70+"]
        df_a["Band"] = pd.cut(df_a["__AGE__"], bins=bins, labels=labels)
        
        grp = df_a.groupby("Band", observed=False)["is_bad"].mean().reset_index()
        grp["Risk"] = grp["is_bad"] * 100
        
        fig = px.bar(grp, x="Band", y="Risk", text=grp["Risk"].round(1),
                     labels={"Band": "Age Group", "Risk": "At Risk (%)"},
                     color_discrete_sequence=[DesignTheme.COLOR_SEQUENCE[0]] * len(grp))
        # Workaround: manually set color
        fig.update_traces(marker_color=DesignTheme.COLOR_SEQUENCE[0])
        
        add_plotly_image(slide, fig, "Age-wise Risk Escalation")

    # ---------------- Slide 6: Young Adult Risk (18-20) ----------------
    if "__AGE__" in view_df.columns:
        grp_ya = view_df[view_df["__AGE__"].between(18,20)]
        if not grp_ya.empty and "Health Status" in grp_ya:
            slide = add_slide(SLIDE_TITLE_ONLY)
            format_title(slide, "Young Adult Risk (18-20 Years)")
            
            # Use 2 Column Layout manually
            # Left: text stats, Right: maybe a placeholder icon or simple stat
            
            any_ab = (grp_ya["Health Status"] != "Healthy").sum()
            txt_lines = [
                f"Sample Size: {len(grp_ya)} individuals",
                f"At Risk: {pct(any_ab, len(grp_ya))}%",
            ]
            
            # Clinical specifics
            if cols_map.get("bmi"):
                b = as_num(grp_ya[cols_map["bmi"]])
                obs = (b>=30).sum()
                txt_lines.append(f"Obesity: {pct(obs, b.notna().sum())}%")
                
            tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = tb.text_frame
            for t in txt_lines:
                p = tf.add_paragraph()
                p.text = "• " + t
                p.font.size = Pt(24)
                p.space_after = Pt(14)
            
    # ---------------- Slide 7: Community-Level Health Burden ----------------
    if cols_map.get("locality") and "Health Status" in view_df.columns:
        slide = add_slide(SLIDE_TITLE_ONLY)
        
        loc = cols_map["locality"]
        df_l = view_df.copy()
        df_l["bad"] = df_l["Health Status"] != "Healthy"
        grp = df_l.groupby(loc, observed=False)["bad"].agg(["count","mean"]).reset_index()
        grp["pct"] = grp["mean"] * 100
        top = grp.sort_values("count", ascending=False).head(10)
        
        fig = px.bar(top, x=loc, y="pct", text=top["pct"].round(1),
                     title="Top Communities by Health Risk",
                     labels={loc: "Community", "pct": "% At Risk"})
        add_plotly_image(slide, fig, "Community Health Burden")

    # ---------------- Slide 8: Community Risk Heatmap ----------------
    if cols_map.get("locality"):
        slide = add_slide(SLIDE_TITLE_ONLY)
        
        loc = cols_map["locality"]
        df_c = view_df.copy()
        df_c["risk"] = (df_c["Health Status"]!="Healthy") if "Health Status" in df_c else 0
        
        # Optional columns
        col_list = ["risk"]
        names_list = ["Overall Risk"]
        
        if cols_map.get("glucose"):
            df_c["diabetes"] = (as_num(df_c[cols_map["glucose"]])>=126)
            col_list.append("diabetes")
            names_list.append("Diabetes")
            
        grp = df_c.groupby(loc, observed=False)[col_list].mean() * 100
        # Sort by first col
        grp = grp.sort_values(col_list[0], ascending=False).head(10)
        grp.columns = names_list
        
        fig = px.imshow(grp, labels=dict(x="Risk Factor", y="Community", color="%"),
                        color_continuous_scale="RdBu_r")
        add_plotly_image(slide, fig, "Community Risk Comparison")

    # ---------------- Slide 9: Review of Top Localities ----------------
    if cols_map.get("locality"):
        slide = add_slide(SLIDE_TITLE_ONLY)
        format_title(slide, "Locality Participation Overview")
        
        loc = cols_map["locality"]
        top_locs = view_df[loc].value_counts().head(5).reset_index()
        top_locs.columns = [loc, "Count"]
        
        fig = px.bar(top_locs, x=loc, y="Count", text="Count",
                     title="Top 5 Localities by Participation",
                     color="Count", color_continuous_scale=DesignTheme.COLOR_SEQUENCE)
        
        add_plotly_image(slide, fig, "Locality Participation")

    # ---------------- Slide 10: Lifestyle Risk Overview ----------------
    # Consolidated Bar Chart
    ls_data = []
    if cols_map.get("tobacco"):
        t_counts = view_df[cols_map["tobacco"]].value_counts(normalize=True).mul(100).head(3)
        for i,v in t_counts.items(): ls_data.append({"Factor": "Tobacco", "Cat": i, "Pct": v})
    if cols_map.get("alcohol"):
        a_counts = view_df[cols_map["alcohol"]].value_counts(normalize=True).mul(100).head(3)
        for i,v in a_counts.items(): ls_data.append({"Factor": "Alcohol", "Cat": i, "Pct": v})
        
    if ls_data:
        slide = add_slide(SLIDE_TITLE_ONLY)
        df_ls = pd.DataFrame(ls_data)
        fig = px.bar(df_ls, x="Cat", y="Pct", color="Factor", title="Lifestyle Factors (%)",
                        text=df_ls["Pct"].round(1), labels={"Cat": "Category","Pct": "Participation (%)"})
        add_plotly_image(slide, fig, "Lifestyle Risk Overview")

    # ---------------- Slide 11: Combined Lifestyle ----------------
    if cols_map.get("tobacco") and cols_map.get("alcohol"):
        slide = add_slide(SLIDE_TITLE_ONLY)
        df_c = view_df.groupby([cols_map["tobacco"], cols_map["alcohol"]], observed=False).size().reset_index(name="Count")
        if not df_c.empty:
            fig = px.treemap(df_c, path=[cols_map["tobacco"], cols_map["alcohol"]], values="Count")
            add_plotly_image(slide, fig, "Alcohol & Tobacco Co-occurrence")

    # ---------------- Slide 12: Socioeconomic ----------------
    soc_col = cols_map.get("income") or cols_map.get("occupation")
    if soc_col and "Health Status" in view_df.columns:
        slide = add_slide(SLIDE_TITLE_ONLY)
        
        grp = view_df.groupby([soc_col, "Health Status"], observed=False).size().reset_index(name="Count")
        # filter top 10 soc
        top_soc = view_df[soc_col].value_counts().head(10).index
        grp = grp[grp[soc_col].isin(top_soc)]
        
        fig = px.bar(grp, x=soc_col, y="Count", color="Health Status", barmode="group")
        add_plotly_image(slide, fig, f"Health Status by {soc_col}")

    # ---------------- Slide 13: Priority Groups ----------------
    slide = add_slide(SLIDE_BULLET)
    format_title(slide, "Priority Groups for Action")
    
    # Textual list
    # Textual list
    priorities = []
    
    # 1. Dual Risk
    if cols_map.get("glucose") and cols_map.get("bp_sys") and cols_map.get("bp_dia"):
        g = as_num(view_df[cols_map["glucose"]])
        s = as_num(view_df[cols_map["bp_sys"]])
        d = as_num(view_df[cols_map["bp_dia"]])
        dual = ((g >= 126) & ((s >= 130) | (d >= 80))).sum()
        if dual > 0:
            priorities.append(f"Individuals with dual risks (Diabetes + Hypertension): {dual:,} identified")

    # 2. Elderly Risk
    if "__AGE__" in view_df.columns and "Health Status" in view_df.columns:
        old = view_df[view_df["__AGE__"] >= 60]
        if len(old) > 0:
            bad_old = (old["Health Status"] != "Healthy").sum()
            old_risk = pct(bad_old, len(old))
            if old_risk > 0:
                priorities.append(f"Elderly population (Age 60+): {old_risk}% showing abnormal health status")

    # 3. Community Risk
    high_risk_locs = []
    if cols_map.get("locality") and "Health Status" in view_df.columns:
        loc = cols_map["locality"]
        grp_l = view_df.groupby(loc, observed=False)["Health Status"].apply(lambda x: (x != "Healthy").mean()).mul(100)
        high_risk_locs = grp_l[grp_l > 40].index.tolist()
        if high_risk_locs:
            count_locs = len(high_risk_locs)
            names = ", ".join(str(x) for x in high_risk_locs[:3])
            priorities.append(f"{count_locs} Communities with >40% risk (e.g., {names})")

    if not priorities:
        priorities.append("No specific high-priority groups identified based on current thresholds.")

    add_styled_bullets(slide, priorities)

    # ---------------- Slide 14: Screening Focus ----------------
    slide = add_slide(SLIDE_BULLET)
    format_title(slide, "Screening & Intervention Focus")
    
    focus_areas = []
    
    # 1. Location based
    if high_risk_locs:
        focus_areas.append("Targeted screening camps in identified high-risk communities")
        
    # 2. Lifestyle
    lifestyle_cols = [c for c in [cols_map.get("tobacco"), cols_map.get("alcohol")] if c]
    if lifestyle_cols:
        focus_areas.append("Lifestyle counseling and cessation programs for tobacco/alcohol users")
        
    # 3. Pre-diabetes (Glucose 100-125)
    if cols_map.get("glucose"):
        g = as_num(view_df[cols_map["glucose"]])
        pred = ((g >= 100) & (g < 126)).sum()
        if pred > 0:
            focus_areas.append(f"Preventive follow-ups for {pred:,} pre-diabetic individuals")
    
    if not focus_areas:
        focus_areas.append("General health awareness and regular screening drives")
        focus_areas.append("Nutritional counseling for the general population")

    add_styled_bullets(slide, focus_areas)

    # ---------------- Slide 15: Summary ----------------
    slide = add_slide(SLIDE_TITLE) # Use title slide for big finish
    
    # Custom BG again
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DesignTheme.DARK_BG
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    p = tb.text_frame.add_paragraph()
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(50)
    p.font.color.rgb = DesignTheme.LIGHT_BG
    p.font.bold = True
    
    p2 = tb.text_frame.add_paragraph()
    p2.text = "Data-Driven Health Insights for a Better Tomorrow"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.color.rgb = DesignTheme.SECONDARY

    # Save
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf